"""Main RAG processor - orchestrates everything."""
from pathlib import Path
import mimetypes
from typing import Dict, List, Optional
import os

from .vector_store import QdrantStore, get_base_dir
from .ocr import OCRProcessor
from .embedder import NomicEmbedder
from .reranker import BGEReranker
from .hybrid_search import HybridSearcher
from .chunking import chunk_text, add_context_to_chunks, generate_summary

class RAGProcessor:
    def __init__(self, qdrant_path: str = None):
        # Resolve absolute path for Qdrant
        if qdrant_path is None:
            data_dir = os.getenv("DATA_DIR")
            base_dir = get_base_dir()
            
            if data_dir:
                if os.path.isabs(data_dir):
                    qdrant_path = os.path.join(data_dir, "qdrant_db")
                else:
                    # Handle relative path by joining with base_dir
                    qdrant_path = str(base_dir / data_dir / "qdrant_db")
            else:
                qdrant_path = str(base_dir / "data" / "qdrant_db")
        
        # Ensure the path is absolute
        if not os.path.isabs(qdrant_path):
            qdrant_path = str(base_dir / qdrant_path)
        
        print(f"Using Qdrant path: {qdrant_path}")
        
        # Initialize models with error handling
        print("Initializing RAG models...")
        
        # Initialize OCR (replacing vision model)
        try:
            self.ocr = OCRProcessor()
            print("OCR processor initialized")
        except Exception as e:
            print(f"Warning: OCR processor failed to initialize: {e}")
            self.ocr = None
        
        # Initialize other components
        try:
            self.embedder = NomicEmbedder()
            print("Embedder initialized")
        except Exception as e:
            print(f"Error: Embedder failed to initialize: {e}")
            raise e  # Embedder is critical
        
        try:
            self.reranker = BGEReranker()
            print("Reranker initialized")
        except Exception as e:
            print(f"Warning: Reranker failed to initialize: {e}")
            self.reranker = None
        
        try:
            self.hybrid = HybridSearcher()
            print("Hybrid searcher initialized")
        except Exception as e:
            print(f"Warning: Hybrid searcher failed to initialize: {e}")
            self.hybrid = None
        
        # Qdrant Vector Store
        try:
            self.vector_store = QdrantStore(path=qdrant_path, collection_name="documents")
            print("Vector store initialized")
        except Exception as e:
            print(f"Error initializing vector store: {e}")
            self.vector_store = None
        
        print("RAG Processor initialized (some components may be disabled).")
    
    def process_file(self, file_path: str, conversation_id: str) -> Dict:
        """Process any file type and store in Qdrant."""
        # Ensure absolute path for file:// URI compatibility
        file_path = Path(file_path).resolve()
        mime_type, _ = mimetypes.guess_type(file_path)
        
        print(f"Processing file: {file_path}, type: {mime_type}")
        
        try:
            # Route to appropriate processor
            text = ""
            method = "text"
            
            if mime_type and mime_type.startswith('image/'):
                text = self._process_image(file_path)
                method = "vision"
            elif file_path.suffix.lower() == '.pdf':
                text = self._process_pdf(file_path)
            elif file_path.suffix.lower() == '.docx':
                text = self._process_docx(file_path)
            elif file_path.suffix.lower() == '.pptx':
                text = self._process_pptx(file_path)
            elif file_path.suffix.lower() in ['.xlsx', '.xls']:
                text = self._process_excel(file_path)
            elif file_path.suffix.lower() in ['.txt', '.md', '.py', '.js', '.ts', '.html', '.css', '.json']:
                text = file_path.read_text(encoding='utf-8', errors='ignore')
            else:
                return {"status": "error", "message": f"Unsupported file type: {file_path.suffix}"}
            
            if not text.strip():
                return {"status": "error", "message": "No text extracted from file"}
            
            # Generate summary for context
            doc_summary = generate_summary(text)
            
            # Chunk with context
            chunks = chunk_text(text, chunk_size=500, overlap=50)
            contextual_chunks = add_context_to_chunks(chunks, doc_summary, file_path.name)
            
            # Embed and store
            print(f"Embedding {len(contextual_chunks)} chunks for {file_path.name}...")
            # Use encode_batch to prevent OOM on large files
            embeddings = self.embedder.encode_batch(contextual_chunks, task_type="search_document", batch_size=32)
            
            ids = [f"{conversation_id}_{file_path.name}_{i}" for i in range(len(contextual_chunks))]
            metadatas = [{
                "conversation_id": conversation_id,
                "file_name": file_path.name,
                "file_type": mime_type or file_path.suffix,
                "chunk_index": i,
                "processing_method": method,
                "parent_summary": doc_summary
            } for i in range(len(contextual_chunks))]
            
            if self.vector_store:
                self.vector_store.add(
                    ids=ids,
                    embeddings=[e.tolist() for e in embeddings],
                    texts=contextual_chunks,
                    metadatas=metadatas
                )
            
            # Index for BM25 (update index for this conversation)
            self._update_bm25_index(conversation_id)
            
            print(f"Successfully processed and stored {file_path.name}")
            
            return {
                "status": "success",
                "chunks": len(chunks),
                "file_name": file_path.name,
                "processing_method": method,
                "extracted_text": text
            }
            
        except Exception as e:
            import traceback
            traceback.print_exc()
            return {"status": "error", "message": str(e)}
    
    def _process_image(self, file_path: Path) -> str:
        """Process image using EasyOCR only."""
        image_path_str = str(file_path)
        print(f"Processing image: {image_path_str}")
        
        # Only use OCR
        if self.ocr is not None:
            try:
                print(f"Attempting OCR extraction for {file_path.name}...")
                ocr_text = self.ocr.extract_text(image_path_str)
                print(f"OCR extraction completed, extracted {len(ocr_text)} characters")
                
                if ocr_text and ocr_text.strip():
                    return ocr_text
                else:
                    print(f"OCR extraction produced empty text for {file_path.name}")
            except Exception as e:
                print(f"OCR extraction failed for {file_path.name}: {e}")
        else:
            print("OCR processor not initialized")
        
        # If OCR fails or is empty, return minimal description
        print(f"Warning: OCR failed or produced no text for image: {file_path.name}")
        
        # Try to provide at least some context based on filename
        filename = file_path.name.lower()
        description_parts = []
        
        if any(ext in filename for ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp']):
            description_parts.append("This is an image file")
        
        # Add basic description based on filename patterns
        if 'screenshot' in filename:
            description_parts.append("It appears to be a screenshot")
        elif 'photo' in filename or 'picture' in filename:
            description_parts.append("It appears to be a photograph")
        elif 'diagram' in filename or 'chart' in filename:
            description_parts.append("It appears to be a diagram or chart")
        elif 'document' in filename or 'scan' in filename:
            description_parts.append("It appears to be a scanned document")
        
        # Add file size info if available
        try:
            file_size = file_path.stat().st_size
            if file_size > 0:
                description_parts.append(f"File size: {file_size:,} bytes")
        except:
            pass
        
        if description_parts:
            return f"[Image Processing Note: Could not extract meaningful text from image. {', '.join(description_parts)}.]"
        else:
            return f"[Image Processing Note: Could not extract meaningful text from image {file_path.name}. The image may not contain readable text or the processing models failed.]"
    
    def _process_pdf(self, file_path: Path) -> str:
        """Extract text and OCR images from PDF using pdfplumber."""
        import pdfplumber
        import tempfile
        
        try:
            text_content = []
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    page_text = [f"### Page {i+1}"]
                    
                    # 1. Extract standard text
                    content = page.extract_text()
                    if content:
                        page_text.append(content)
                    
                    # 2. Extract tables (pdfplumber specialty)
                    tables = page.extract_tables()
                    for table in tables:
                        if table:
                            table_str = "\n| " + " | ".join([str(c).replace("\n", " ") if c else "" for c in table[0]]) + " |\n"
                            table_str += "| " + " | ".join(["---" for _ in table[0]]) + " |\n"
                            for row in table[1:]:
                                table_str += "| " + " | ".join([str(c).replace("\n", " ") if c else "" for c in row]) + " |\n"
                            page_text.append(table_str)
                    
                    # 3. OCR Images on page
                    for j, image in enumerate(page.images):
                        try:
                            # Extract image bytes
                            # pdfplumber provides image object with stream
                            # But often easier to use page.to_image().original
                            # For simplicity and robustness, we use a crop approach if needed
                            # or just detect that an image exists and try to OCR the whole page as image if text is low
                            pass # Placeholder for advanced per-image OCR
                        except:
                            pass
                    
                    # If very little text was extracted, try OCR on the whole page
                    if not content or len(content.strip()) < 50:
                        try:
                            # Render page to image for OCR
                            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                                page.to_image(resolution=200).save(tmp.name)
                                tmp_path = Path(tmp.name)
                            
                            ocr_result = self._process_image(tmp_path)
                            if ocr_result and "Image Processing Note" not in ocr_result and ocr_result.strip():
                                page_text.append(f"\n[OCR Text from Page Image:\n{ocr_result}\n]")
                            
                            if os.path.exists(tmp_path):
                                os.unlink(tmp_path)
                        except Exception as e:
                            print(f"Failed to OCR PDF page {i+1}: {e}")

                    text_content.append("\n".join(page_text))
            
            return "\n\n---\n\n".join(text_content)
        except Exception as e:
            print(f"PDF processing failed: {e}")
            # Fallback to simple pypdf if pdfplumber fails
            try:
                import pypdf
                with open(file_path, 'rb') as f:
                    reader = pypdf.PdfReader(f)
                    return "\n\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
            except:
                return f"Error processing PDF: {str(e)}"
    
    def _process_docx(self, file_path: Path) -> str:
        """Extract text from DOCX."""
        import docx
        doc = docx.Document(file_path)
        return "\n\n".join([para.text for para in doc.paragraphs])

    def _process_pptx(self, file_path: Path) -> str:
        """Robustly extract text from PPTX, including placeholders, groups, tables, and notes."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        import tempfile
        
        def get_shape_text(shape):
            """Deeply extract text from shapes and subgroups."""
            text = ""
            # Handle standard text frames
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = "".join(run.text for run in paragraph.runs).strip()
                    if para_text:
                        text += para_text + "\n"
            # Handle tables
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                for row in shape.table.rows:
                    text += "| " + " | ".join([cell.text_frame.text.replace("\n", " ").strip() for cell in row.cells]) + " |\n"
            # Recursive group traversal
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for s in shape.shapes:
                    text += get_shape_text(s) + "\n"
            # General text attribute fallback
            elif hasattr(shape, "text") and shape.text:
                text += shape.text.strip() + "\n"
            return text.strip()

        try:
            prs = Presentation(file_path)
            all_slides_content = []
            
            print(f"Starting robust extraction for PPTX: {file_path.name}")
            
            for i, slide in enumerate(prs.slides):
                slide_parts = [f"--- SLIDE {i+1} ---"]
                
                # 1. Capture Slide Title specifically
                if slide.shapes.title:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        slide_parts.append(f"TITLE: {title_text}")
                
                # 2. Extract from all shapes using robust recursive method
                for shape in slide.shapes:
                    # Skip title as we already got it
                    if shape == slide.shapes.title:
                        continue
                        
                    shape_text = get_shape_text(shape)
                    if shape_text:
                        # Prefix content with slide context to survive chunking
                        slide_context = f"[Slide {i+1}"
                        if slide.shapes.title and slide.shapes.title.text.strip():
                             slide_context += f": {slide.shapes.title.text.strip()[:20]}..."
                        slide_context += "] "
                        
                        slide_parts.append(slide_context + shape_text)
                    
                    # 3. Image OCR extraction for embedded pictures
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            with tempfile.NamedTemporaryFile(suffix=f".{image.ext}", delete=False) as tmp:
                                tmp.write(image.blob)
                                tmp_path = Path(tmp.name)
                            
                            ocr_text = self._process_image(tmp_path)
                            # Only include if we found actual text (avoiding the generic fallback note)
                            if ocr_text and "[Image Processing Note:" not in ocr_text and ocr_text.strip():
                                slide_parts.append(f"[TEXT DETECTED IN IMAGE: {ocr_text.strip()}]")
                            
                            if os.path.exists(tmp_path): os.unlink(tmp_path)
                        except Exception as e:
                            print(f"Slide {i+1} picture OCR failed: {e}")

                # 4. Capture Speaker Notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_parts.append(f"SPEAKER NOTES: {notes}")

                slide_content = "\n".join(slide_parts)
                all_slides_content.append(slide_content)
                print(f"PPTX Slide {i+1} extracted: {len(slide_content)} chars")

            final_text = "\n\n".join(all_slides_content)
            print(f"PPTX {file_path.name} total extraction: {len(final_text)} characters")
            return final_text
            
        except Exception as e:
            print(f"PPTX processing failed: {e}")
            import traceback
            traceback.print_exc()
            return f"[Error processing PPTX {file_path.name}: {str(e)}]"
        except Exception as e:
            print(f"PPTX processing failed: {e}")
            return f"Error processing PPTX: {str(e)}"
    
    def _process_excel(self, file_path: Path) -> str:
        """Extract text from Excel as context-rich rows using fast libraries (pylightxl/xlrd) with openpyxl fallback."""
        text = []
        
        try:
            # 1. Try pylightxl for .xlsx (Fastest, Low Memory)
            if file_path.suffix.lower() == '.xlsx':
                try:
                    import pylightxl as xl
                    db = xl.readxl(fn=str(file_path))
                    for sheet in db.ws_names:
                        text.append(f"### Sheet: {sheet}")
                        # pylightxl rows are generators
                        rows = list(db.ws(ws=sheet).rows)
                        if not rows: continue
                        
                        # Context-aware extraction
                        headers = [str(cell).strip() if cell else f"Col_{j+1}" for j, cell in enumerate(rows[0])]
                        
                        for r_idx, row in enumerate(rows[1:], start=2):
                            row_parts = []
                            has_data = False
                            for c_idx, cell in enumerate(row):
                                if cell and str(cell).strip():
                                    header = headers[c_idx] if c_idx < len(headers) else f"Col_{c_idx+1}"
                                    val = str(cell).replace("\n", " ").strip()
                                    row_parts.append(f"[{header}: {val}]")
                                    has_data = True
                            if has_data:
                                text.append(f"Row {r_idx}: " + " ".join(row_parts))
                        text.append("\n")
                    return "\n".join(text)
                except ImportError:
                    pass # pylightxl not installed
                except Exception as e:
                    print(f"pylightxl failed: {e}, falling back...")

            # 2. Try xlrd for .xls (Legacy support)
            if file_path.suffix.lower() == '.xls':
                try:
                    import xlrd
                    book = xlrd.open_workbook(str(file_path))
                    for sheet in book.sheets():
                        text.append(f"### Sheet: {sheet.name}")
                        if sheet.nrows == 0: continue
                        
                        headers = [str(cell.value).strip() if cell.value else f"Col_{j+1}" for j, cell in enumerate(sheet.row(0))]
                        
                        for r_idx in range(1, sheet.nrows):
                            row = sheet.row(r_idx)
                            row_parts = []
                            has_data = False
                            for c_idx, cell in enumerate(row):
                                if cell.value and str(cell.value).strip():
                                    header = headers[c_idx] if c_idx < len(headers) else f"Col_{c_idx+1}"
                                    val = str(cell.value).replace("\n", " ").strip()
                                    row_parts.append(f"[{header}: {val}]")
                                    has_data = True
                            if has_data:
                                text.append(f"Row {r_idx+1}: " + " ".join(row_parts))
                        text.append("\n")
                    return "\n".join(text)
                except ImportError:
                    pass # xlrd not installed
                except Exception as e:
                    print(f"xlrd failed: {e}, falling back...")

            # 3. Fallback to openpyxl (Slow but reliable, supports .xlsm etc)
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=True)
            
            for sheet in wb.worksheets:
                text.append(f"### Sheet: {sheet.title}")
                rows_iter = sheet.iter_rows(values_only=True)
                try:
                    header_row = next(rows_iter)
                except StopIteration:
                    continue
                    
                headers = [str(cell).strip() if cell is not None else f"Col_{j+1}" for j, cell in enumerate(header_row)]
                
                for r_idx, row in enumerate(rows_iter, start=2):
                    row_parts = []
                    has_data = False
                    for c_idx, cell in enumerate(row):
                        if cell is not None and str(cell).strip():
                            header = headers[c_idx] if c_idx < len(headers) else f"Col_{c_idx+1}"
                            val = str(cell).replace("\n", " ").strip()
                            row_parts.append(f"[{header}: {val}]")
                            has_data = True
                    if has_data:
                        text.append(f"Row {r_idx}: " + " ".join(row_parts))
                text.append("\n")
                
            return "\n".join(text)

        except Exception as e:
            return f"[Excel Processing Error: {str(e)}]"
    
    def _update_bm25_index(self, conversation_id: str):
        """Update BM25 index with conversation documents."""
        if not self.vector_store:
            return
            
        documents = self.vector_store.get_all_for_conversation(conversation_id)
        if documents:
            self.hybrid.index_corpus(documents)
    
    def retrieve_context(self, query: str, conversation_id: str, top_k: int = 3) -> str:
        """Retrieve relevant context using hybrid search + reranking."""
        try:
            # Update BM25 index for this conversation if not already loaded or if different
            # For simplicity, we just reload it. Optimization: Check if conversation_id changed.
            self._update_bm25_index(conversation_id)
            
            # Dense retrieval
            # Use raw query without heuristic expansion for better precision
            query_embedding = self.embedder.encode([query], task_type="search_query")[0]
            
            dense_results = []
            if self.vector_store:
                # Search for documents
                results = self.vector_store.query(
                    query_embedding=query_embedding.tolist(),
                    conversation_id=conversation_id,
                    top_k=15
                )
                dense_results = results
            
            dense_docs = [r["content"] for r in dense_results]
            
            # Sparse retrieval (if hybrid searcher is available)
            sparse_docs = []
            if self.hybrid is not None:
                try:
                    # Use raw query for BM25 as well
                    sparse_docs = self.hybrid.search(query, top_k=15)
                except Exception as e:
                    print(f"Sparse retrieval failed: {e}")
            
            # Merge with RRF if we have both dense and sparse results
            merged = []
            if dense_docs and sparse_docs and self.hybrid is not None:
                try:
                    merged = self.hybrid.merge_results(dense_docs, sparse_docs)
                except Exception as e:
                    print(f"Merge failed: {e}")
                    merged = dense_docs[:15]  # Fallback to dense results
            elif dense_docs:
                merged = dense_docs[:15]
            elif sparse_docs:
                merged = sparse_docs[:15]
            
            # Rerank if reranker is available
            if merged and self.reranker is not None:
                try:
                    reranked = self.reranker.rerank(query, merged[:20], top_k=top_k)
                    context = "\n\n---\n\n".join([doc for score, doc in reranked])
                    return context
                except Exception as e:
                    print(f"Reranking failed: {e}")
                    # Fallback to simple selection
                    context = "\n\n---\n\n".join(merged[:top_k])
                    return context
            elif merged:
                # No reranker, just take top results
                context = "\n\n---\n\n".join(merged[:top_k])
                return context
            
            return ""
        except Exception as e:
            print(f"Retrieval failed: {e}")
            return ""
    
    def get_file_content(self, file_names: List[str], conversation_id: str) -> str:
        """
        Retrieve full content of specific files from Qdrant.
        Used for explicit file attachments in chat.
        If not found in DB, falls back to reading from uploads directory.
        """
        try:
            if not file_names:
                return ""
                
            print(f"Retrieving content for files: {file_names} in conv: {conversation_id}")
            
            # 1. Try to get from Vector Store
            found_files = {}
            if self.vector_store:
                results = self.vector_store.get_files(file_names, conversation_id)
                for r in results:
                    fname = r["metadata"].get("file_name", "")
                    if fname:
                        if fname not in found_files:
                            found_files[fname] = []
                        found_files[fname].append((r["metadata"].get("chunk_index", 0), r["content"]))
            
            # 2. Process each requested file
            context_parts = []
            files_with_errors = []
            files_not_found = []
            files_empty = []
            
            for fname in file_names:
                content = ""
                source = "vector_db"
                has_error = False
                
                if fname in found_files:
                    # Sort chunks by index
                    chunks = sorted(found_files[fname], key=lambda x: x[0])
                    content = "".join([c[1] for c in chunks])
                    
                    # Check if content contains error messages
                    if any(error_indicator in content for error_indicator in [
                        "[Error reading file",
                        "[Excel Processing Error",
                        "[Excel Processing Note",
                        "[Image Processing Note",
                        "[Error processing PDF",
                        "[Error processing PPTX"
                    ]):
                        has_error = True
                        files_with_errors.append(fname)
                else:
                    # Fallback to uploads folder
                    source = "upload"
                    media_dir = os.getenv("MEDIA_DIR")
                    if media_dir:
                        # If path is relative, anchor it to project root
                        if not os.path.isabs(media_dir):
                            base_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../"))
                            uploads_dir = os.path.join(base_dir, media_dir, "uploads")
                        else:
                            uploads_dir = os.path.join(media_dir, "uploads")
                    else:
                        uploads_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../uploads"))
                    
                    file_path = os.path.join(uploads_dir, fname)
                    
                    if os.path.exists(file_path):
                        print(f"File {fname} not found in VectorDB, reading from uploads: {file_path}")
                        try:
                            # Use helper methods to extract content
                            fpath = Path(file_path)
                            mime_type, _ = mimetypes.guess_type(fpath)
                            
                            if mime_type and mime_type.startswith('image/'):
                                content = self._process_image(fpath)
                            elif fpath.suffix.lower() == '.pdf':
                                content = self._process_pdf(fpath)
                            elif fpath.suffix.lower() == '.docx':
                                content = self._process_docx(fpath)
                            elif fpath.suffix.lower() in ['.xlsx', '.xls']:
                                content = self._process_excel(fpath)
                            else:
                                # Text-based
                                content = fpath.read_text(encoding='utf-8', errors='ignore')
                            
                            # Check for error messages in extracted content
                            if any(error_indicator in content for error_indicator in [
                                "[Error reading file",
                                "[Excel Processing Error",
                                "[Excel Processing Note",
                                "[Image Processing Note",
                                "[Error processing PDF",
                                "[Error processing PPTX"
                            ]):
                                has_error = True
                                files_with_errors.append(fname)
                            elif not content.strip():
                                files_empty.append(fname)
                                content = f"[File Processing Note: File '{fname}' is empty or contains no extractable text]"
                                
                        except Exception as e:
                            print(f"Failed to read file from uploads {fname}: {e}")
                            content = f"[File Processing Error: Cannot read file '{fname}'. Error: {str(e)}]"
                            has_error = True
                            files_with_errors.append(fname)
                    else:
                        print(f"File {fname} not found in VectorDB OR uploads folder.")
                        files_not_found.append(fname)
                        continue # Skip missing files
                
                if content:
                    # Add error/warning header if needed
                    if has_error:
                        context_parts.append(f"\n--- Content of Attached File: {fname} ({source}) [ERROR READING FILE] ---\n")
                    elif fname in files_empty:
                        context_parts.append(f"\n--- Content of Attached File: {fname} ({source}) [EMPTY FILE] ---\n")
                    else:
                        context_parts.append(f"\n--- Content of Attached File: {fname} ({source}) ---\n")
                    context_parts.append(content)
            
            # Add summary of file status if there are issues
            if files_with_errors or files_not_found or files_empty:
                summary_parts = ["\n--- FILE PROCESSING SUMMARY ---"]
                if files_with_errors:
                    summary_parts.append(f"Files with errors (cannot be read): {', '.join(files_with_errors)}")
                if files_not_found:
                    summary_parts.append(f"Files not found: {', '.join(files_not_found)}")
                if files_empty:
                    summary_parts.append(f"Empty files (no extractable content): {', '.join(files_empty)}")
                summary_parts.append("The assistant should acknowledge these issues when answering questions about these files.")
                context_parts.append("\n".join(summary_parts))
            
            if not context_parts:
                print(f"No content found for files: {file_names}")
                return f"[File Processing Note: No content found for any of the requested files: {', '.join(file_names)}]"
            
            return "\n".join(context_parts)
            
        except Exception as e:
            print(f"Failed to get file content: {e}")
            import traceback
            traceback.print_exc()
            return f"[File Processing Error: System error retrieving file content. Error: {str(e)}]"

    def search(self, query: str, conversation_id: str, top_k: int = 10) -> List[Dict]:
        """
        Structured search returning list of results with metadata.
        Used by semantic search endpoints.
        """
        try:
            if not self.vector_store:
                return []

            # Dense retrieval only for now to ensure metadata preservation
            query_embedding = self.embedder.encode([query], task_type="search_query")[0]
            
            results = self.vector_store.query(
                query_embedding=query_embedding.tolist(),
                conversation_id=conversation_id,
                top_k=top_k
            )
            
            formatted_results = []
            for r in results:
                formatted_results.append({
                    "content": r["content"],
                    "metadata": r["metadata"],
                    "similarity": r["score"]
                })
            
            return formatted_results
        except Exception as e:
            print(f"Search failed: {e}")
            return []
